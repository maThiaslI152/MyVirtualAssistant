from typing import Dict, Any, List, Optional, Union, Set, Tuple
import logging
import os
from pathlib import Path
import PyPDF2
import pytesseract
from PIL import Image
import cv2
import numpy as np
from langchain.text_splitter import RecursiveCharacterTextSplitter
import chardet
import magic
from dataclasses import dataclass
import asyncio
import docx
import openpyxl
import pandas as pd
from git import Repo
import hashlib
from datetime import datetime
import json
import ast
import javalang
import esprima
import pptx
import striprtf
import hglib
import difflib
import radon.complexity as cc
import radon.metrics as metrics
from radon.visitors import ComplexityVisitor
from radon.raw import analyze
import lizard
import mccabe
import tempfile
from typing_extensions import TypedDict
import git

class VersionInfo(TypedDict):
    history: List[Dict[str, Any]]
    current_hash: str
    last_modified: str
    total_commits: int
    vcs_type: str

class CodeMetrics(TypedDict):
    complexity: Dict[str, int]
    maintainability: Dict[str, float]
    raw_metrics: Dict[str, int]
    halstead_metrics: Dict[str, float]
    cognitive_complexity: int

@dataclass
class ProcessedFile:
    content: str
    metadata: Dict[str, Any]
    chunks: List[Dict[str, Any]]
    file_type: str
    mime_type: str
    version_info: Optional[VersionInfo] = None
    code_metrics: Optional[CodeMetrics] = None
    diff_info: Optional[Dict[str, Any]] = None

class FileProcessorService:
    def __init__(
        self,
        base_path: str,
        git_repo_path: Optional[str] = None,
        hg_repo_path: Optional[str] = None,
        version_history_limit: int = 10
    ):
        self.base_path = base_path
        self.git_repo = None
        self.hg_client = None
        self.version_history_limit = version_history_limit
        self.logger = logging.getLogger(__name__)

        # Initialize Git client if repo path provided
        if git_repo_path:
            try:
                self.git_repo = git.Repo(git_repo_path)
            except Exception as e:
                self.logger.warning(f"Could not initialize Git client: {str(e)}")

        # Initialize Mercurial client if repo path provided
        if hg_repo_path:
            try:
                self.hg_client = hglib.open(hg_repo_path)
            except Exception as e:
                self.logger.warning(f"Could not initialize Mercurial client: {str(e)}")

        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len,
            separators=["\n\n", "\n", " ", ""]
        )

    def _detect_encoding(self, file_path: str) -> str:
        """Detect file encoding."""
        with open(file_path, 'rb') as f:
            result = chardet.detect(f.read())
        return result['encoding'] or 'utf-8'

    def _get_file_metadata(self, file_path: str) -> Dict[str, Any]:
        """Get file metadata."""
        file_stat = os.stat(file_path)
        mime_type = magic.from_file(file_path, mime=True)
        return {
            'filename': os.path.basename(file_path),
            'file_size': file_stat.st_size,
            'created_time': file_stat.st_ctime,
            'modified_time': file_stat.st_mtime,
            'mime_type': mime_type,
            'extension': os.path.splitext(file_path)[1].lower()
        }

    def _process_pdf(self, file_path: str) -> ProcessedFile:
        """Process PDF files."""
        try:
            content = []
            metadata = self._get_file_metadata(file_path)
            metadata['pages'] = 0
            metadata['pdf_metadata'] = {}

            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                metadata['pages'] = len(pdf_reader.pages)
                metadata['pdf_metadata'] = pdf_reader.metadata

                for page_num in range(len(pdf_reader.pages)):
                    page = pdf_reader.pages[page_num]
                    content.append(f"Page {page_num + 1}:\n{page.extract_text()}")

            full_content = "\n\n".join(content)
            chunks = self._create_chunks(full_content, metadata)

            return ProcessedFile(
                content=full_content,
                metadata=metadata,
                chunks=chunks,
                file_type='pdf',
                mime_type='application/pdf'
            )

        except Exception as e:
            self.logger.error(f"Error processing PDF {file_path}: {str(e)}")
            raise

    def _process_image(self, file_path: str) -> ProcessedFile:
        """Process image files including screenshots."""
        try:
            metadata = self._get_file_metadata(file_path)
            
            # Read image
            image = cv2.imread(file_path)
            if image is None:
                raise ValueError(f"Could not read image: {file_path}")

            # Resize if too large
            height, width = image.shape[:2]
            if height > (1920, 1080)[1] or width > (1920, 1080)[0]:
                scale = min((1920, 1080)[0] / width, (1920, 1080)[1] / height)
                image = cv2.resize(image, None, fx=scale, fy=scale)

            # Preprocess image for better OCR
            enhanced_image = self._enhance_image_for_ocr(image)
            
            # Perform OCR
            content = pytesseract.image_to_string(enhanced_image, lang='+'.join(['eng']))
            
            # Get OCR confidence data
            ocr_data = pytesseract.image_to_data(enhanced_image, output_type=pytesseract.Output.DICT)
            ocr_confidence = ocr_data.get('conf', [])
            
            # Extract additional metadata
            metadata.update({
                'image_size': f"{width}x{height}",
                'channels': image.shape[2] if len(image.shape) > 2 else 1,
                'ocr_confidence': ocr_confidence
            })

            chunks = self._create_chunks(content, metadata)

            return ProcessedFile(
                content=content,
                metadata=metadata,
                chunks=chunks,
                file_type='image',
                mime_type=metadata['mime_type']
            )

        except Exception as e:
            self.logger.error(f"Error processing image {file_path}: {str(e)}")
            raise

    def _process_code_file(self, file_path: str) -> ProcessedFile:
        """Process code files with syntax highlighting and structure analysis."""
        try:
            metadata = self._get_file_metadata(file_path)
            extension = metadata['extension']

            # Read file with proper encoding
            encoding = self._detect_encoding(file_path)
            with open(file_path, 'r', encoding=encoding) as f:
                content = f.read()

            # Extract code structure
            structure = self._analyze_code_structure(content, extension)
            metadata['code_structure'] = structure

            # Create chunks with context
            chunks = self._create_code_chunks(content, structure, metadata)

            return ProcessedFile(
                content=content,
                metadata=metadata,
                chunks=chunks,
                file_type='code',
                mime_type=metadata['mime_type']
            )

        except Exception as e:
            self.logger.error(f"Error processing code file {file_path}: {str(e)}")
            raise

    def _enhance_image_for_ocr(self, image: np.ndarray) -> np.ndarray:
        """Enhance image for better OCR results."""
        # Convert to grayscale
        gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
        
        # Apply adaptive thresholding
        thresh = cv2.adaptiveThreshold(
            gray, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C, cv2.THRESH_BINARY, 11, 2
        )
        
        # Denoise
        denoised = cv2.fastNlMeansDenoising(thresh)
        
        # Enhance contrast
        clahe = cv2.createCLAHE(clipLimit=2.0, tileGridSize=(8,8))
        enhanced = clahe.apply(denoised)
        
        # Deskew
        coords = np.column_stack(np.where(enhanced > 0))
        angle = cv2.minAreaRect(coords)[-1]
        if angle < -45:
            angle = 90 + angle
        (h, w) = enhanced.shape[:2]
        center = (w // 2, h // 2)
        M = cv2.getRotationMatrix2D(center, angle, 1.0)
        rotated = cv2.warpAffine(
            enhanced, M, (w, h), flags=cv2.INTER_CUBIC, borderMode=cv2.BORDER_REPLICATE
        )
        
        return rotated

    def _process_word_document(self, file_path: str) -> ProcessedFile:
        """Process Word documents."""
        try:
            doc = docx.Document(file_path)
            content = []
            metadata = self._get_file_metadata(file_path)
            
            # Extract text from paragraphs
            for para in doc.paragraphs:
                content.append(para.text)
            
            # Extract tables
            for table in doc.tables:
                for row in table.rows:
                    content.append(' | '.join(cell.text for cell in row.cells))
            
            # Extract metadata
            metadata.update({
                'core_properties': {
                    'author': doc.core_properties.author,
                    'created': doc.core_properties.created,
                    'modified': doc.core_properties.modified,
                    'title': doc.core_properties.title,
                    'subject': doc.core_properties.subject,
                    'keywords': doc.core_properties.keywords
                },
                'sections': len(doc.sections),
                'tables': len(doc.tables),
                'paragraphs': len(doc.paragraphs)
            })
            
            full_content = "\n\n".join(content)
            chunks = self._create_chunks(full_content, metadata)
            
            return ProcessedFile(
                content=full_content,
                metadata=metadata,
                chunks=chunks,
                file_type='word',
                mime_type='application/vnd.openxmlformats-officedocument.wordprocessingml.document'
            )
            
        except Exception as e:
            self.logger.error(f"Error processing Word document {file_path}: {str(e)}")
            raise

    def _process_excel_file(self, file_path: str) -> ProcessedFile:
        """Process Excel files."""
        try:
            wb = openpyxl.load_workbook(file_path, data_only=True)
            content = []
            metadata = self._get_file_metadata(file_path)
            
            # Process each sheet
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                content.append(f"Sheet: {sheet_name}")
                
                # Get sheet data
                data = []
                for row in sheet.iter_rows(values_only=True):
                    data.append([str(cell) if cell is not None else '' for cell in row])
                
                # Convert to DataFrame for better processing
                df = pd.DataFrame(data)
                content.append(df.to_string())
            
            # Extract metadata
            metadata.update({
                'sheets': wb.sheetnames,
                'active_sheet': wb.active.title,
                'properties': {
                    'creator': wb.properties.creator,
                    'created': wb.properties.created,
                    'modified': wb.properties.modified,
                    'title': wb.properties.title,
                    'subject': wb.properties.subject,
                    'keywords': wb.properties.keywords
                }
            })
            
            full_content = "\n\n".join(content)
            chunks = self._create_chunks(full_content, metadata)
            
            return ProcessedFile(
                content=full_content,
                metadata=metadata,
                chunks=chunks,
                file_type='excel',
                mime_type='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
            )
            
        except Exception as e:
            self.logger.error(f"Error processing Excel file {file_path}: {str(e)}")
            raise

    def _process_powerpoint(self, file_path: str) -> ProcessedFile:
        """Process PowerPoint presentations."""
        try:
            prs = pptx.Presentation(file_path)
            content = []
            metadata = self._get_file_metadata(file_path)
            
            # Extract text from slides
            for slide in prs.slides:
                slide_content = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_content.append(shape.text)
                content.append(f"Slide {slide.slide_id}:\n" + "\n".join(slide_content))
            
            # Extract metadata
            metadata.update({
                'slides': len(prs.slides),
                'core_properties': {
                    'author': prs.core_properties.author,
                    'created': prs.core_properties.created,
                    'modified': prs.core_properties.modified,
                    'title': prs.core_properties.title,
                    'subject': prs.core_properties.subject,
                    'keywords': prs.core_properties.keywords
                }
            })
            
            full_content = "\n\n".join(content)
            chunks = self._create_chunks(full_content, metadata)
            
            return ProcessedFile(
                content=full_content,
                metadata=metadata,
                chunks=chunks,
                file_type='powerpoint',
                mime_type='application/vnd.openxmlformats-officedocument.presentationml.presentation'
            )
            
        except Exception as e:
            self.logger.error(f"Error processing PowerPoint file {file_path}: {str(e)}")
            raise

    def _process_rtf(self, file_path: str) -> ProcessedFile:
        """Process RTF files."""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                rtf_content = file.read()
            
            # Convert RTF to plain text
            content = striprtf.rtf_to_text(rtf_content)
            metadata = self._get_file_metadata(file_path)
            
            chunks = self._create_chunks(content, metadata)
            
            return ProcessedFile(
                content=content,
                metadata=metadata,
                chunks=chunks,
                file_type='rtf',
                mime_type='application/rtf'
            )
            
        except Exception as e:
            self.logger.error(f"Error processing RTF file {file_path}: {str(e)}")
            raise

    def _analyze_code_structure(self, content: str, extension: str) -> Dict[str, Any]:
        """Enhanced code structure analysis."""
        structure = {
            'imports': [],
            'classes': [],
            'functions': [],
            'variables': [],
            'comments': [],
            'dependencies': set(),
            'complexity': {
                'cyclomatic': 0,
                'cognitive': 0
            },
            'metrics': {
                'lines': len(content.splitlines()),
                'non_empty_lines': len([l for l in content.splitlines() if l.strip()]),
                'comment_lines': 0
            }
        }

        try:
            if extension == '.py':
                # Python specific analysis
                tree = ast.parse(content)
                
                # Analyze imports
                for node in ast.walk(tree):
                    if isinstance(node, ast.Import):
                        for name in node.names:
                            structure['imports'].append(f"import {name.name}")
                            structure['dependencies'].add(name.name)
                    elif isinstance(node, ast.ImportFrom):
                        structure['imports'].append(f"from {node.module} import {', '.join(n.name for n in node.names)}")
                        structure['dependencies'].add(node.module)
                    elif isinstance(node, ast.ClassDef):
                        structure['classes'].append({
                            'name': node.name,
                            'bases': [b.id for b in node.bases if isinstance(b, ast.Name)],
                            'methods': [m.name for m in node.body if isinstance(m, ast.FunctionDef)],
                            'decorators': [d.id for d in node.decorator_list if isinstance(d, ast.Name)]
                        })
                    elif isinstance(node, ast.FunctionDef):
                        structure['functions'].append({
                            'name': node.name,
                            'args': [arg.arg for arg in node.args.args],
                            'returns': ast.unparse(node.returns) if node.returns else None,
                            'decorators': [d.id for d in node.decorator_list if isinstance(d, ast.Name)]
                        })
                        # Calculate cyclomatic complexity
                        structure['complexity']['cyclomatic'] += len([n for n in ast.walk(node) if isinstance(n, (ast.If, ast.While, ast.For, ast.Try, ast.ExceptHandler))])
                    elif isinstance(node, ast.Assign):
                        for target in node.targets:
                            if isinstance(target, ast.Name):
                                structure['variables'].append({
                                    'name': target.id,
                                    'value': ast.unparse(node.value)
                                })
                    elif isinstance(node, ast.Expr) and isinstance(node.value, ast.Str):
                        structure['comments'].append(node.value.s)
                        structure['metrics']['comment_lines'] += 1

            elif extension in ['.js', '.ts']:
                # JavaScript/TypeScript specific analysis
                tree = esprima.parseScript(content, {'loc': True, 'range': True})
                
                for node in tree.body:
                    if node.type == 'ImportDeclaration':
                        structure['imports'].append(f"import {node.source.value}")
                        structure['dependencies'].add(node.source.value)
                    elif node.type == 'ClassDeclaration':
                        structure['classes'].append({
                            'name': node.id.name,
                            'methods': [m.key.name for m in node.body.body if m.type == 'MethodDefinition'],
                            'extends': node.superClass.name if node.superClass else None
                        })
                    elif node.type == 'FunctionDeclaration':
                        structure['functions'].append({
                            'name': node.id.name,
                            'params': [p.name for p in node.params],
                            'is_async': getattr(node, 'async', False),
                            'is_generator': getattr(node, 'generator', False)
                        })
                    elif node.type == 'VariableDeclaration':
                        for decl in node.declarations:
                            structure['variables'].append({
                                'name': decl.id.name,
                                'value': esprima.unparse(decl.init) if decl.init else None
                            })
                    elif node.type == 'ExpressionStatement' and node.expression.type == 'CallExpression':
                        if node.expression.callee.name == 'require':
                            structure['dependencies'].add(node.expression.arguments[0].value)

            elif extension == '.java':
                # Java specific analysis
                tree = javalang.parse.parse(content)
                
                for type_decl in tree.types:
                    if isinstance(type_decl, javalang.tree.ClassDeclaration):
                        structure['classes'].append({
                            'name': type_decl.name,
                            'extends': type_decl.extends.name if type_decl.extends else None,
                            'implements': [i.name for i in type_decl.implements] if type_decl.implements else [],
                            'methods': [m.name for m in type_decl.methods],
                            'fields': [f.name for f in type_decl.fields]
                        })
                    elif isinstance(type_decl, javalang.tree.InterfaceDeclaration):
                        structure['classes'].append({
                            'name': type_decl.name,
                            'extends': [e.name for e in type_decl.extends] if type_decl.extends else [],
                            'methods': [m.name for m in type_decl.methods]
                        })

                for imp in tree.imports:
                    structure['imports'].append(f"import {imp.path}")
                    structure['dependencies'].add(imp.path)

        except Exception as e:
            self.logger.warning(f"Error analyzing code structure: {str(e)}")

        # Convert sets to lists for JSON serialization
        structure['dependencies'] = list(structure['dependencies'])
        return structure

    def _create_code_chunks(
        self,
        content: str,
        structure: Dict[str, Any],
        metadata: Dict[str, Any]
    ) -> List[Dict[str, Any]]:
        """Create chunks from code with context preservation."""
        chunks = []
        lines = content.split('\n')
        current_chunk = []
        current_context = []

        for i, line in enumerate(lines):
            current_chunk.append(line)
            
            # Check if line is part of structure
            for category, patterns in structure.items():
                if any(line.strip() in pattern for pattern in patterns):
                    current_context.append(f"{category}: {line.strip()}")

            # Create chunk if size threshold reached
            if len('\n'.join(current_chunk)) >= 1000:
                chunk_content = '\n'.join(current_chunk)
                chunks.append({
                    'content': chunk_content,
                    'metadata': {
                        **metadata,
                        'context': current_context[-3:],  # Keep last 3 context items
                        'line_start': i - len(current_chunk) + 1,
                        'line_end': i
                    }
                })
                # Keep overlap
                overlap_start = max(0, len(current_chunk) - 200)
                current_chunk = current_chunk[overlap_start:]
                current_context = current_context[-3:]

        # Add remaining content
        if current_chunk:
            chunks.append({
                'content': '\n'.join(current_chunk),
                'metadata': {
                    **metadata,
                    'context': current_context,
                    'line_start': len(lines) - len(current_chunk),
                    'line_end': len(lines) - 1
                }
            })

        return chunks

    def _create_chunks(self, content: str, metadata: Dict[str, Any]) -> List[Dict[str, Any]]:
        """Create chunks from content with metadata."""
        chunks = []
        for chunk in self.text_splitter.split_text(content):
            chunks.append({
                'content': chunk,
                'metadata': metadata
            })
        return chunks

    def _calculate_code_metrics(self, content: str, extension: str) -> CodeMetrics:
        """Calculate comprehensive code metrics."""
        metrics_data = {
            'complexity': {},
            'maintainability': {},
            'raw_metrics': {},
            'halstead_metrics': {},
            'cognitive_complexity': 0
        }

        try:
            # Calculate cyclomatic complexity
            if extension == '.py':
                visitor = ComplexityVisitor.from_code(content)
                metrics_data['complexity']['cyclomatic'] = visitor.complexity
            else:
                metrics_data['complexity']['cyclomatic'] = cc.cc_visit(content)[0].complexity

            # Calculate maintainability index
            raw_metrics = analyze(content)
            metrics_data['raw_metrics'] = {
                'loc': raw_metrics.loc,
                'lloc': raw_metrics.lloc,
                'sloc': raw_metrics.sloc,
                'comments': raw_metrics.comments,
                'multi': raw_metrics.multi,
                'blank': raw_metrics.blank
            }
            
            mi = metrics.mi_visit(content, multi=True)
            metrics_data['maintainability'] = {
                'index': mi,
                'rank': metrics.mi_rank(mi)
            }

            # Calculate Halstead metrics
            halstead = metrics.h_visit(content)
            if halstead:
                metrics_data['halstead_metrics'] = {
                    'vocabulary': halstead[0].vocabulary,
                    'length': halstead[0].length,
                    'volume': halstead[0].volume,
                    'difficulty': halstead[0].difficulty,
                    'effort': halstead[0].effort,
                    'time': halstead[0].time,
                    'bugs': halstead[0].bugs
                }

            # Calculate cognitive complexity
            if extension == '.py':
                with tempfile.NamedTemporaryFile(suffix='.py', mode='w', delete=False) as temp_file:
                    temp_file.write(content)
                    temp_file.flush()
                    analysis = lizard.analyze_file.analyze_source_code(temp_file.name, content)
                    metrics_data['cognitive_complexity'] = analysis.average_cyclomatic_complexity
                os.unlink(temp_file.name)
            else:
                ast_tree = ast.parse(content)
                visitor = mccabe.PathGraphingAstVisitor()
                visitor.visitFunctionDef(ast_tree)
                metrics_data['cognitive_complexity'] = visitor.complexity

        except Exception as e:
            self.logger.warning(f"Error calculating code metrics: {str(e)}")

        return metrics_data

    def get_version_history(self, file_path: str) -> List[Dict[str, Any]]:
        """Get version history for a file."""
        version_info = []
        
        # Try Git first
        if self.git_repo:
            try:
                # Get relative path from git repo root
                rel_path = os.path.relpath(file_path, self.git_repo.working_dir)
                commits = list(self.git_repo.iter_commits(paths=rel_path, max_count=self.version_history_limit))
                
                for commit in commits:
                    version_info.append({
                        'version': commit.hexsha,
                        'author': commit.author.name,
                        'date': commit.committed_datetime.isoformat(),
                        'message': commit.message.strip(),
                        'vcs_type': 'git'
                    })
            except Exception as e:
                self.logger.warning(f"Error getting Git version info: {str(e)}")
        
        # Try Mercurial if Git failed
        if not version_info and self.hg_client:
            try:
                # Get relative path from hg repo root
                rel_path = os.path.relpath(file_path, self.hg_client.root())
                log = self.hg_client.log(rel_path, limit=self.version_history_limit)
                
                for entry in log:
                    version_info.append({
                        'version': entry[1].decode('utf-8'),
                        'author': entry[5].decode('utf-8'),
                        'date': datetime.fromtimestamp(entry[4]).isoformat(),
                        'message': entry[6].decode('utf-8'),
                        'vcs_type': 'hg'
                    })
            except Exception as e:
                self.logger.warning(f"Error getting Mercurial version info: {str(e)}")
        
        return version_info

    def get_file_diff(self, file_path: str, old_version: str, new_version: str) -> str:
        """Get diff between two versions of a file."""
        if self.git_repo:
            try:
                # Get relative path from git repo root
                rel_path = os.path.relpath(file_path, self.git_repo.working_dir)
                old_content = self.git_repo.git.show(f"{old_version}:{rel_path}")
                new_content = self.git_repo.git.show(f"{new_version}:{rel_path}")
                return self._compute_diff(old_content, new_content)
            except Exception as e:
                self.logger.warning(f"Error getting Git diff: {str(e)}")
        
        elif self.hg_client:
            try:
                # Get relative path from hg repo root
                rel_path = os.path.relpath(file_path, self.hg_client.root())
                old_content = self.hg_client.cat([rel_path], rev=old_version).decode('utf-8')
                new_content = self.hg_client.cat([rel_path], rev=new_version).decode('utf-8')
                return self._compute_diff(old_content, new_content)
            except Exception as e:
                self.logger.warning(f"Error getting Mercurial diff: {str(e)}")
        
        return ""

    async def process_file(self, file_path: str) -> ProcessedFile:
        """Process a file based on its type."""
        try:
            file_path = Path(file_path)
            if not file_path.exists():
                raise FileNotFoundError(f"File not found: {file_path}")

            mime_type = magic.from_file(str(file_path), mime=True)
            extension = file_path.suffix.lower()

            # Get version information
            version_info = self.get_version_history(str(file_path))

            # Get diff information if version info exists
            diff_info = None
            if version_info and len(version_info) > 1:
                diff_info = self.get_file_diff(
                    str(file_path),
                    version_info[1]['version'],
                    version_info[0]['version']
                )

            # Get code metrics for code files
            code_metrics = None
            if extension in ['.py', '.js', '.ts', '.java', '.cpp', '.c', '.h', '.hpp', '.cs', '.go', '.rb', '.php', '.swift', '.kt', '.rs', '.scala', '.html', '.css', '.scss', '.sass', '.less', '.xml', '.json', '.yaml', '.yml', '.toml', '.ini', '.md', '.rst', '.tex', '.sql', '.sh', '.bash', '.ps1', '.bat', '.cmd']:
                with open(file_path, 'r', encoding=self._detect_encoding(str(file_path))) as f:
                    content = f.read()
                code_metrics = self._calculate_code_metrics(content, extension)

            if mime_type == 'application/pdf':
                return self._process_pdf(str(file_path))
            elif mime_type.startswith('image/'):
                return self._process_image(str(file_path))
            elif extension in ['.py', '.js', '.ts', '.java', '.cpp', '.c', '.h', '.hpp', '.cs', '.go', '.rb', '.php', '.swift', '.kt', '.rs', '.scala', '.html', '.css', '.scss', '.sass', '.less', '.xml', '.json', '.yaml', '.yml', '.toml', '.ini', '.md', '.rst', '.tex', '.sql', '.sh', '.bash', '.ps1', '.bat', '.cmd']:
                return self._process_code_file(str(file_path))
            elif extension == '.docx':
                return self._process_word_document(str(file_path))
            elif extension in ['.xlsx', '.xls']:
                return self._process_excel_file(str(file_path))
            elif extension == '.pptx':
                return self._process_powerpoint(str(file_path))
            elif extension == '.rtf':
                return self._process_rtf(str(file_path))
            else:
                raise ValueError(f"Unsupported file type: {mime_type}")

        except Exception as e:
            self.logger.error(f"Error processing file {file_path}: {str(e)}")
            raise

    async def process_files(self, file_paths: List[str]) -> List[ProcessedFile]:
        """Process multiple files concurrently."""
        async def process_single_file(file_path: str) -> ProcessedFile:
            return await self.process_file(file_path)

        tasks = [process_single_file(file_path) for file_path in file_paths]
        return await asyncio.gather(*tasks, return_exceptions=True) 